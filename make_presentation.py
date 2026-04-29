"""Generates the thesis presentation PPTX and speaker script."""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT   = Path(__file__).resolve().parent
IMGS   = {
    "l1": ROOT / "rl-coverage/level1_decoder/curves.png",
    "l2": ROOT / "rl-coverage/level2_cpu_196bin/curves_cpu.png",
    "l3": ROOT / "rl-coverage/level3_chains/curves_chains.png",
    "l4": ROOT / "rl-coverage/level4_shadow/curves_l6.png",
    "l5": ROOT / "rl-coverage/level5_real_rtl/curves_l5.png",
    "l6": ROOT / "rl-coverage/level6_rvc/rvc_comparison.png",
    "l7": ROOT / "rl-coverage/level7_stimulus/l7_comparison.png",
}
OUT_PPTX   = ROOT / "ibex_coverage_presentation.pptx"
OUT_SCRIPT = ROOT / "speaker_script.md"

# ── light theme colours ────────────────────────────────────────────────────
C_BG      = RGBColor(0xFF, 0xFF, 0xFF)   # white
C_BG2     = RGBColor(0xF4, 0xF7, 0xFB)  # very light blue-gray (title slide bg)
C_ACCENT  = RGBColor(0x1D, 0x6F, 0xC5)  # professional blue
C_DARK    = RGBColor(0x1A, 0x1A, 0x2E)  # near-black for headings
C_BODY    = RGBColor(0x33, 0x33, 0x33)  # dark gray for body
C_MUTED   = RGBColor(0x66, 0x77, 0x88)  # muted for subtitles
C_BOX_BG  = RGBColor(0xEB, 0xF3, 0xFC)  # light blue for boxes
C_GREEN   = RGBColor(0x1A, 0x7A, 0x3C)
C_ORANGE  = RGBColor(0xC0, 0x5A, 0x00)

W, H = Inches(13.33), Inches(7.5)


# ── helpers ────────────────────────────────────────────────────────────────
def new_prs():
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    return prs


def blank(prs, bg=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill  = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = bg or C_BG
    return slide


def txt(slide, text, x, y, w, h, size=20, bold=False, color=C_BODY,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = align
    r  = p.add_run()
    r.text = text
    r.font.size   = Pt(size)
    r.font.bold   = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def bar(slide, x, y, w, h=0.045, color=C_ACCENT):
    sh = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Emu(int(h * 914400)))
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background()


def box(slide, x, y, w, h, fill=C_BOX_BG, border=C_ACCENT):
    sh = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = border
    sh.line.width = Emu(12700)
    return sh


def add_title(slide, title, subtitle=None):
    txt(slide, title, 0.55, 0.22, 12.2, 0.9, size=34, bold=True, color=C_DARK)
    bar(slide, 0.55, 1.17, 12.2)
    if subtitle:
        txt(slide, subtitle, 0.55, 1.28, 12.2, 0.55, size=16,
            color=C_MUTED, italic=True)


def img(slide, path, x, y, w):
    if Path(path).exists():
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))


def badge(slide, label, value, x, y, w=2.7, h=1.2):
    box(slide, x, y, w, h, fill=C_BOX_BG, border=C_ACCENT)
    txt(slide, value, x+0.1, y+0.08, w-0.2, 0.65, size=24, bold=True,
        color=C_ACCENT, align=PP_ALIGN.CENTER)
    txt(slide, label, x+0.1, y+0.72, w-0.2, 0.42, size=12,
        color=C_MUTED, align=PP_ALIGN.CENTER)


def bullet_list(slide, items, x, y, w, size=19, gap=0.52):
    """items: str  →  bullet  |  (heading, body) → bold heading + indented body"""
    cy = y
    for item in items:
        if isinstance(item, tuple):
            h, b = item
            txt(slide, h, x, cy, w, gap, size=size, bold=True, color=C_ACCENT)
            cy += gap - 0.05
            txt(slide, b, x+0.2, cy, w-0.2, gap+0.1, size=size-2, color=C_BODY)
            cy += gap + 0.1
        else:
            txt(slide, f"•  {item}", x, cy, w, gap+0.05, size=size, color=C_BODY)
            cy += gap


# ══════════════════════════════════════════════════════════════════════════
#  SLIDES
# ══════════════════════════════════════════════════════════════════════════
def build(prs):

    # ── 1. TITLE ──────────────────────────────────────────────────────────
    s = blank(prs, bg=C_BG2)
    bar(s, 0, 0, 13.33, h=0.35, color=C_ACCENT)
    txt(s, "Automated Hardware Verification",
        1.0, 1.5, 11.3, 1.1, size=42, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
    txt(s, "with Reinforcement Learning",
        1.0, 2.55, 11.3, 0.95, size=38, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    bar(s, 3.6, 3.6, 6.13)
    txt(s, "Coverage Experiments on the Ibex RISC-V Core",
        1.0, 3.75, 11.3, 0.65, size=20, color=C_MUTED, align=PP_ALIGN.CENTER)
    bar(s, 0, 7.15, 13.33, h=0.35, color=C_ACCENT)

    # ── 2. WHAT IS HARDWARE VERIFICATION ──────────────────────────────────
    s = blank(prs)
    add_title(s, "Hardware Verification", "Why does it matter?")
    bullet_list(s, [
        "A CPU has tens of thousands of internal wires, registers, and FSM states",
        "Every state must be tested before the chip is manufactured — bugs in silicon are permanent",
        "Standard approach: thousands of hand-written test programs + commercial simulators",
        "lowRISC verifies Ibex with 1530 manual tests → 88.7% branch coverage",
        "This thesis: can an automated agent reach a comparable number with free, open-source tools?",
    ], 0.55, 1.85, 12.2, size=21, gap=0.88)

    # ── 3. TYPES OF COVERAGE ──────────────────────────────────────────────
    s = blank(prs)
    add_title(s, "Types of Coverage")

    types = [
        ("Toggle",     "Each wire toggled  0→1  and  1→0",                   "structural, auto"),
        ("Branch",     "Each if/case branch taken at least once",             "structural, auto"),
        ("Line/Block", "Each RTL line / code block executed",                 "structural, auto"),
        ("Expression", "Each boolean sub-expression evaluated 0 and 1",      "structural, auto"),
        ("FSM",        "Each state and transition in a state machine visited","structural, auto"),
        ("Functional", "Engineer-defined scenarios: hazards, sequences, corners", "manual"),
    ]
    for i, (name, desc, kind) in enumerate(types):
        col = 0 if i < 3 else 6.65
        row = i % 3
        bx, by = col + 0.55, 1.85 + row * 1.55
        box(s, bx, by, 6.0, 1.35,
            fill=C_BOX_BG if kind == "structural, auto" else RGBColor(0xFF,0xF4,0xE8),
            border=C_ACCENT if kind == "structural, auto" else C_ORANGE)
        txt(s, name, bx+0.15, by+0.08, 5.7, 0.5, size=18, bold=True, color=C_DARK)
        txt(s, desc, bx+0.15, by+0.52, 5.7, 0.55, size=14, color=C_BODY)
        txt(s, kind, bx+0.15, by+1.0, 5.7, 0.3, size=12, italic=True, color=C_MUTED)

    txt(s, "This project:  toggle (L1–L7)  +  functional (L9)",
        0.55, 7.05, 12.2, 0.38, size=16, bold=True,
        color=C_ACCENT, align=PP_ALIGN.CENTER)

    # ── 4. TYPES OF INSTRUCTIONS ──────────────────────────────────────────
    s = blank(prs)
    add_title(s, "RISC-V Instruction Types", "Context for the coverage experiments")

    itypes = [
        ("R-type",      "Register–register  (ADD, SUB, AND, MUL, SLL …)"),
        ("I-type",      "Immediate operand  (ADDI, LW, LH, LB, JALR …)"),
        ("S-type",      "Stores  (SW, SH, SB)"),
        ("B-type",      "Branches  (BEQ, BNE, BLT, BGE …)"),
        ("U-type",      "Upper immediate  (LUI, AUIPC)"),
        ("J-type",      "Jumps  (JAL)"),
        ("CSR",         "Control/Status Registers  (mcycle, minstret, mcause …)"),
        ("RVC",         "Compressed 16-bit variants of the above  (C.ADD, C.LW …)"),
    ]
    cols = [(0.55, 1.85), (6.8, 1.85)]
    for i, (name, desc) in enumerate(itypes):
        cx, cy = cols[i // 4]
        ry = cy + (i % 4) * 1.28
        box(s, cx, ry, 5.95, 1.15, fill=C_BOX_BG, border=C_ACCENT)
        txt(s, name, cx+0.15, ry+0.08, 5.6, 0.45, size=17, bold=True, color=C_ACCENT)
        txt(s, desc, cx+0.15, ry+0.52, 5.6, 0.5,  size=14, color=C_BODY)

    txt(s, "Levels 1–7 progressively add more types → wider coverage",
        0.55, 7.05, 12.2, 0.38, size=16, bold=True,
        color=C_MUTED, align=PP_ALIGN.CENTER)

    # ── 5. IBEX + TECH STACK ──────────────────────────────────────────────
    s = blank(prs)
    add_title(s, "Ibex & Tech Stack")

    # left: ibex
    txt(s, "Target: Ibex RISC-V Core", 0.55, 1.85, 6.0, 0.5, size=20,
        bold=True, color=C_DARK)
    bullet_list(s, [
        "~15,000 lines of SystemVerilog",
        "Used in Google's OpenTitan silicon",
        "Minimal config: no cache, no PMP, no debug",
        "Public baseline: 88.7% (lowRISC, 1530 tests)",
    ], 0.55, 2.45, 6.0, size=18, gap=0.58)

    bar(s, 6.85, 1.85, 0.04, h=5.2, color=RGBColor(0xCC,0xCC,0xCC))

    # right: tools
    txt(s, "Tools", 7.1, 1.85, 6.0, 0.5, size=20, bold=True, color=C_DARK)
    tools = [
        ("Verilator",  "Compiles SV RTL → C++ simulator  (free)"),
        ("cocotb",     "Python ↔ Simulator bridge"),
        ("RVFI",       "Observes every instruction retiring"),
        ("Gymnasium",  "RL environment wrapper"),
        ("SB3 / PPO",  "Proximal Policy Optimization agent"),
        ("Shadow env", "Pure-Python coverage model  ×1500 faster"),
    ]
    cy = 2.45
    for name, desc in tools:
        txt(s, name, 7.1, cy, 2.0, 0.42, size=17, bold=True, color=C_ACCENT)
        txt(s, desc, 9.2, cy, 4.0, 0.42, size=16, color=C_BODY)
        cy += 0.6

    # ── 6. THE LOOP ───────────────────────────────────────────────────────
    s = blank(prs)
    add_title(s, "The Generate → Simulate → Measure Loop")

    steps = [
        ("1  Generate",  "Agent picks\ninstructions"),
        ("2  Simulate",  "Verilator runs\nthe program on Ibex"),
        ("3  Measure",   "Parse\ncoverage.dat"),
        ("4  Reward",    "New bins hit\n→ reward signal"),
    ]
    for i, (label, sub) in enumerate(steps):
        bx = 0.6 + i * 3.1
        box(s, bx, 2.1, 2.7, 1.4, fill=C_BOX_BG, border=C_ACCENT)
        txt(s, label, bx+0.1, 2.2, 2.5, 0.55, size=18, bold=True,
            color=C_ACCENT, align=PP_ALIGN.CENTER)
        txt(s, sub, bx+0.1, 2.75, 2.5, 0.65, size=14,
            color=C_BODY, align=PP_ALIGN.CENTER)
        if i < 3:
            txt(s, "→", bx+2.75, 2.55, 0.35, 0.6, size=26,
                bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

    txt(s, "Shadow env: replace Verilator with pure Python  →  ~1500× faster  →  enables RL training",
        0.55, 3.85, 12.2, 0.55, size=18, color=C_MUTED, align=PP_ALIGN.CENTER)
    txt(s, "Real RTL: each episode ~2–3 s  →  budget dozens to hundreds of episodes",
        0.55, 4.45, 12.2, 0.55, size=18, color=C_MUTED, align=PP_ALIGN.CENTER)

    bar(s, 0.55, 5.2, 12.2, color=RGBColor(0xCC,0xCC,0xCC))

    txt(s, "Levels 1–4: shadow only       |       Levels 5–9: real Verilator RTL",
        0.55, 5.35, 12.2, 0.55, size=17, bold=True,
        color=C_DARK, align=PP_ALIGN.CENTER)

    # ── CHART SLIDES (label + full-width chart) ───────────────────────────
    chart_slides = [
        ("L1 — Combinational Decoder",
         "PPO on a 2107-bin shadow.  Result: ~100% reachable bins.",
         "l1"),
        ("L2 — CPU Shadow  (196 bins)",
         "LLM4DV benchmark.  PPO hits 196/196.  Shadow → real RTL: 196/196 match.",
         "l2"),
        ("L3 — Chained Coverage  (1739 bins)",
         "Multi-cycle dependencies.  PPO: 99%+.  Random lags for the first time.",
         "l3"),
        ("L4 — Full Shadow  (5615 bins)  —  PPO 100× faster than random",
         "Shadow-to-real validation: 425/425 and 1499/1499 bins match.",
         "l4"),
        ("L5 — First Real Verilator Toggle Coverage",
         "All three PPO variants plateau at 56.2% — same as random.  Ceiling = encoder, not algorithm.",
         "l5"),
        ("L6 — +RVC Compressed Instructions",
         "compressed_decoder: 0% → 97%.  Cumulative gain: +1.3 pp.  PPO still = random.",
         "l6"),
        ("L7 — Stimulus Engineering vs Algorithm Tuning",
         "4 env changes (mem prepop, trap handler, 29 CSRs, AUIPC) → +10 pp.  "
         "L7 PPO (300 eps): 67.58%  vs  L7 random (30 eps): 66.27%.",
         "l7"),
    ]

    for title, subtitle, key in chart_slides:
        s = blank(prs)
        add_title(s, title, subtitle)
        if IMGS[key].exists():
            img(s, IMGS[key], 0.55, 1.75, 12.2)

    # ── L9 ────────────────────────────────────────────────────────────────
    s = blank(prs)
    add_title(s, "L9 — Rich Functional Coverage  (~15,217 bins)",
              "Shifting from toggle RTL to engineer-defined functional bins")
    bullet_list(s, [
        ("Coverage model",
         "9 categories: SEEN, OPERAND_PAIR, RAW_HAZARD, SEQUENCE, CORNER_CASES …"),
        ("Python shadow",
         "Full RISC-V ISA simulator with real control flow.  "
         "Critical fix: sequential execution vs. PC tracking  (shadow reported 3× more bins before fix)."),
        ("Speed",
         "~2 s/episode shadow  vs  ~30 s Verilator  →  15× faster  →  RL training practical."),
        ("Baseline to beat",
         "Constrained random (480 steps, 300 episodes)  →  45.10%"),
    ], 0.55, 1.85, 12.2, size=19, gap=0.92)

    # ── CONCLUSIONS ───────────────────────────────────────────────────────
    s = blank(prs)
    add_title(s, "Conclusions & Contributions")
    bullet_list(s, [
        "5615-bin Python shadow for RV32I/M — validated bin-by-bin, zero divergence from real RTL",
        "Reachable ceiling characterised at 75.9% for minimal Ibex",
        "Main finding: stimulus engineering moves coverage ~10× more than algorithm tuning",
        "L7 random (30 eps): 66.27%  vs  L5 PPO (300 eps): 56.20%  — 10 pp gap",
        "L7 PPO (300 eps): 67.58%  — same ceiling as random, small RL advantage (+1.3 pp)",
        "Open-source infrastructure: encoders, parsers, trap scaffolding, unreachability classifier",
    ], 0.55, 1.85, 12.2, size=20, gap=0.85)

    # ── Q&A ───────────────────────────────────────────────────────────────
    s = blank(prs, bg=C_BG2)
    bar(s, 0, 0, 13.33, h=0.35, color=C_ACCENT)
    txt(s, "Thank you", 1.0, 2.3, 11.3, 1.2, size=52, bold=True,
        color=C_DARK, align=PP_ALIGN.CENTER)
    bar(s, 3.6, 3.65, 6.13)
    txt(s, "Questions?", 1.0, 3.85, 11.3, 0.9, size=32,
        color=C_MUTED, align=PP_ALIGN.CENTER)
    bar(s, 0, 7.15, 13.33, h=0.35, color=C_ACCENT)


# ══════════════════════════════════════════════════════════════════════════
#  SPEAKER SCRIPT
# ══════════════════════════════════════════════════════════════════════════
SCRIPT = """# Speaker Script

---

## Slide 1 — Title
Good morning. Today I'm presenting my thesis on automated hardware verification using reinforcement learning, with experiments on the Ibex RISC-V processor.

---

## Slide 2 — Hardware Verification
Hardware verification is one of the most expensive parts of chip design. A modern CPU has an enormous number of internal states and you need to make sure every relevant one is exercised before the chip goes to manufacturing — bugs in silicon are permanent and costly. The industry standard is thousands of hand-written tests on commercial simulators. lowRISC verifies Ibex with 1530 manual tests and reaches 88.7% coverage. The question of this thesis is: can an automated agent approach that number using only free tools?

---

## Slide 3 — Types of Coverage
There are several ways to measure how well you've exercised a design. Toggle, branch, line, and expression coverage are structural — Verilator extracts them automatically from the RTL without any manual effort. FSM coverage tracks state machine transitions. Functional coverage is different: you define it yourself, specifying the scenarios that matter for correctness — instruction combinations, data hazards, corner cases. This project uses toggle coverage in Levels 1–7, and moves to rich functional coverage in Level 9.

---

## Slide 4 — RISC-V Instruction Types
RISC-V instructions come in several formats. R-type operates register to register. I-type has an immediate operand and includes loads. S-type is stores. B-type is branches. U-type is upper immediate — LUI and AUIPC. J-type is jumps. CSR instructions read and write control/status registers like the cycle counter or cause register. RVC is the compressed 16-bit variant of the above. Each level of the experiments progressively adds more of these types to the generator, which is directly what drives coverage up.

---

## Slide 5 — Ibex & Tech Stack
Ibex is a real open-source RISC-V processor — about 15,000 lines of SystemVerilog, used in Google's OpenTitan chip. Verilator compiles that RTL into a fast C++ simulator. Cocotb bridges Python to the simulator so we can write programs and feed them to the CPU from Python. RVFI is the formal interface that exposes every instruction retiring from the pipeline. Gymnasium and PPO from Stable-Baselines3 provide the RL framework. The shadow environment is a pure-Python replica of the coverage logic — about 1500 times faster than Verilator, which is what makes RL training practical.

---

## Slide 6 — The Loop
The fundamental loop is: generate a program, simulate it on Ibex, measure coverage, reward the agent for new bins hit, repeat. Against real Verilator each episode takes 2–3 seconds. Against the shadow it takes milliseconds. Levels 1–4 train exclusively on the shadow. Levels 5–9 run against real RTL.

---

## Slide 7 — L1 Chart
The first experiment targets the Ibex instruction decoder — purely combinational, no clock, no state. PPO trained on a Python shadow saturates almost all reachable bins: 2041 of 2107. The remaining 66 are ISA-unreachable. The loop works.

---

## Slide 8 — L2 Chart
Level 2 uses the LLM4DV benchmark — 196 bins on the full CPU. The best published result using Claude 3.5 Sonnet was 5.61%. PPO hits 196 of 196, and the same program run on real Verilator Ibex also hits 196 of 196. Shadow-to-real fidelity is confirmed.

---

## Slide 9 — L3 Chart
Level 3 extends to 1739 bins with chained multi-cycle dependencies. PPO hits 99%+. Random lags for the first time — sequential planning starts to matter.

---

## Slide 10 — L4 Chart
Level 4 builds a comprehensive 5615-bin shadow for all of RV32I/M. PPO saturates it in about 2 minutes; random takes 53 minutes. That's the 100× speedup. Shadow-to-real validation shows zero divergence: 425 of 425 bins match on one agent, 1499 of 1499 on another.

---

## Slide 11 — L5 Chart
First run on real Verilator toggle coverage. Surprising result: all three PPO variants plateau at exactly 56.2% — same as random, even at 300 episodes. The per-module breakdown explains why: compressed decoder at 0%, load/store unit at 38%, CSR registers at 10%. The ceiling was the encoder, not the algorithm.

---

## Slide 12 — L6 Chart
Adding 16 RVC compressed instructions fixes the compressed decoder — it jumps from 0% to 97% in one run. But the cumulative gain is only 1.3 percentage points because that module is small. PPO still equals random. The unreachability analysis reveals the real gap: the exception path requires ECALL/EBREAK and a trap handler.

---

## Slide 13 — L7 Chart
This is the main finding. Four changes to the environment, none to the algorithm: memory prepopulated with a diverse pattern instead of a constant — load/store unit 38% to 92%. A trap handler so ECALL and EBREAK return safely — entire exception path unlocked. 29 CSR addresses instead of 5. AUIPC added. Result: random in 30 episodes hits 66.27% — 10 points above the best PPO at 300 episodes on the old environment. PPO on the same L7 environment reaches 67.58% at 300 episodes — a small advantage over random, both hitting the same ceiling. The ceiling is set by what you can generate, not by the optimizer.

---

## Slide 14 — L9
Level 9 shifts the target to rich functional coverage — about 15,000 bins covering instruction sequences, operand pairs, RAW hazards, and corner cases. A full RISC-V ISA simulator was built in Python with correct control flow. A critical bug was found and fixed: the initial version executed instructions sequentially ignoring branches, reporting 3× more coverage than Verilator. After the fix the shadow-to-Verilator ratio stabilised at about 0.84. The constrained random baseline is 45.10%. PPO training is ongoing.

---

## Slide 15 — Conclusions
The main contributions: a 5615-bin Python shadow with zero RTL divergence, a characterised reachable ceiling of 75.9%, and the core finding that stimulus engineering dominates algorithm tuning by roughly a factor of 10 on this class of problem. The L7 PPO experiment confirms that once the stimulus space is rich enough, RL holds a small advantage over random but both are bounded by the same ceiling.

---

## Slide 16 — Q&A
Thank you. Happy to take questions.
"""


def main():
    prs = new_prs()
    build(prs)
    prs.save(str(OUT_PPTX))
    print(f"Saved  {OUT_PPTX}")
    OUT_SCRIPT.write_text(SCRIPT, encoding="utf-8")
    print(f"Saved  {OUT_SCRIPT}")


if __name__ == "__main__":
    main()
